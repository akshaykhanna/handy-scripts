from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont

def create_bullet_point_image(font_path, bullet_char="•", font_size=20, font_color="black"):
    # Create a blank image for bullet point
    bullet_image = Image.new("RGB", (font_size, font_size), color="white")
    draw = ImageDraw.Draw(bullet_image)

    # Load font
    font = ImageFont.truetype(font_path, size=font_size)

    # Draw bullet point character
    draw.text((0, 0), bullet_char, font=font, fill=font_color)

    return bullet_image

def create_slide_image(title, bullet_points, font_family, image_width=720, image_height=405):
    # Create a blank image
    image = Image.new("RGB", (image_width, image_height), color="white")
    draw = ImageDraw.Draw(image)

    # Set font properties
    title_font = ImageFont.load_default()  # Default font for title
    bullet_font = ImageFont.load_default()  # Default font for bullet points

    # Draw title
    title_width, title_height = draw.textsize(title)
    draw.text(((image_width - title_width) / 2, 20), title, font=title_font, fill="black")

    # Draw bullet points
    y = 60
    bullet_point_image = create_bullet_point_image(font_family)
    for bullet_point in bullet_points:
        image.paste(bullet_point_image, (50, y))
        draw.text((70, y), bullet_point, font=bullet_font, fill="black")
        y += 30

    return image

def save_image(image, filename):
    image.save(filename)

def create_ppt_slide(title, bullet_points, font_family, pptx_filename):
    prs = Presentation()

    slide_layout = prs.slide_layouts[5]  # Use the layout for title and content

    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = "\n".join(["\u2022 " + bp for bp in bullet_points])

    prs.save(pptx_filename)

# Example usage
title = "Sample Slide"
bullet_points = ["Bullet Point 1", "Bullet Point 2", "Bullet Point 3"]
font_family = "Arial"  # Specify the font family name of your choice

# Create and save the image
slide_image = create_slide_image(title, bullet_points, font_family)
save_image(slide_image, "slide_image.png")

# # Create and save the PowerPoint slide
# create_ppt_slide(title, bullet_points, font_family, "slide.pptx")
